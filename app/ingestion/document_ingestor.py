from __future__ import annotations

import hashlib
import logging
from pathlib import Path
import shutil
import uuid

from app.chunking.chunker import TextChunker
from app.config.settings import settings
from app.indexing.store import ArtifactStore
from app.ingestion.pdf_ingestor import PDFIngestor
from app.ingestion.video_transcriber import TranscriptSegment, VideoTranscriber
from app.schemas.models import DocumentRecord, PageRecord, path_to_str
from app.utils.text import detect_language

logger = logging.getLogger(__name__)

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".m4v", ".mkv", ".webm"}
SUPPORTED_EXTENSIONS = SUPPORTED_EXTENSIONS | VIDEO_EXTENSIONS


class DocumentIngestor:
    def __init__(self) -> None:
        self.pdf_ingestor = PDFIngestor()
        self.store = ArtifactStore()
        self.chunker = TextChunker()
        self.video_transcriber = VideoTranscriber()

    def ingest(
        self,
        source_path: Path,
        course_id: str,
        title: str | None = None,
        uploader_teacher_id: str | None = None,
        source_filename: str | None = None,
    ) -> DocumentRecord:
        ext = source_path.suffix.lower()
        if ext == ".pdf":
            return self.pdf_ingestor.ingest_pdf(
                source_path=source_path,
                course_id=course_id,
                title=title,
                uploader_teacher_id=uploader_teacher_id,
                source_filename=source_filename,
            )
        if ext == ".docx":
            return self._ingest_docx(
                source_path=source_path,
                course_id=course_id,
                title=title,
                uploader_teacher_id=uploader_teacher_id,
                source_filename=source_filename,
            )
        if ext == ".pptx":
            return self._ingest_pptx(
                source_path=source_path,
                course_id=course_id,
                title=title,
                uploader_teacher_id=uploader_teacher_id,
                source_filename=source_filename,
            )
        if ext in VIDEO_EXTENSIONS:
            return self.ingest_video(
                source_path=source_path,
                course_id=course_id,
                title=title,
                uploader_teacher_id=uploader_teacher_id,
                source_filename=source_filename,
            )
        raise ValueError(f"Unsupported file format: {ext}. Supported: .pdf, .docx, .pptx, .mp4, .mov, .m4v, .mkv, .webm")

    def prepare_video_upload(
        self,
        source_path: Path,
        course_id: str,
        title: str | None = None,
        uploader_teacher_id: str | None = None,
        source_filename: str | None = None,
    ) -> DocumentRecord:
        ext = source_path.suffix.lower()
        if ext not in VIDEO_EXTENSIONS:
            raise ValueError(f"Unsupported video format: {ext}")
        document_id, copied = self._copy_source(source_path, course_id, ext)
        doc_title = title or source_path.stem
        document = DocumentRecord(
            document_id=document_id,
            course_id=course_id,
            document_title=doc_title,
            source_pdf=path_to_str(copied.resolve()),
            page_count=0,
            mime_type=_mime_for_ext(ext),
            status="uploaded",
        )
        checksum = hashlib.sha256(copied.read_bytes()).hexdigest()
        self.store.create_document(
            document,
            uploader_teacher_id=uploader_teacher_id,
            source_filename=source_filename or source_path.name,
            checksum_sha256=checksum,
            mime_type=document.mime_type,
        )
        self.store.update_document_status(document.document_id, status="transcribing")
        document.status = "transcribing"
        return document

    def process_video_document(self, document_id: str, course_id: str) -> DocumentRecord:
        document = self.store.get_document(document_id)
        if document is None:
            raise ValueError(f"Document not found: {document_id}")
        source_path = Path(document.source_pdf)
        if not source_path.exists():
            raise ValueError(f"Video file missing on disk: {source_path}")
        segments = self.video_transcriber.transcribe(source_path)
        pages = self._build_video_pages(
            course_id=course_id,
            document_id=document_id,
            document_title=document.document_title,
            segments=segments,
        )
        self.store.create_pages(pages)
        chunks = self.chunker.chunk_pages(pages)
        self.store.upsert_chunks_for_document(course_id=course_id, document_id=document_id, chunks=chunks)
        self.store.update_document_status(document_id, status="ingested", page_count=len(pages))
        return DocumentRecord(
            document_id=document.document_id,
            course_id=document.course_id,
            document_title=document.document_title,
            source_pdf=document.source_pdf,
            page_count=len(pages),
            mime_type=getattr(document, "mime_type", _mime_for_ext(source_path.suffix.lower())),
            status="ingested",
        )

    def ingest_video(
        self,
        source_path: Path,
        course_id: str,
        title: str | None = None,
        uploader_teacher_id: str | None = None,
        source_filename: str | None = None,
    ) -> DocumentRecord:
        doc = self.prepare_video_upload(
            source_path=source_path,
            course_id=course_id,
            title=title,
            uploader_teacher_id=uploader_teacher_id,
            source_filename=source_filename,
        )
        return self.process_video_document(doc.document_id, course_id=course_id)

    def _ingest_docx(
        self,
        source_path: Path,
        course_id: str,
        title: str | None,
        uploader_teacher_id: str | None,
        source_filename: str | None,
    ) -> DocumentRecord:
        try:
            from docx import Document as DocxDocument
        except Exception as exc:  # pragma: no cover - dependency may be missing
            raise RuntimeError("python-docx is not installed. Add python-docx to requirements.") from exc

        document_id, copied = self._copy_source(source_path, course_id, ".docx")
        doc_title = title or source_path.stem
        payload = DocxDocument(copied)
        units: list[str] = []

        for para in payload.paragraphs:
            text = (para.text or "").strip()
            if text:
                units.append(text)

        for table in payload.tables:
            for row in table.rows:
                row_text = " | ".join((cell.text or "").strip() for cell in row.cells if (cell.text or "").strip())
                if row_text:
                    units.append(row_text)

        pages = self._build_text_units_as_pages(
            course_id=course_id,
            document_id=document_id,
            document_title=doc_title,
            units=units,
            unit_prefix="block",
        )
        checksum = hashlib.sha256(copied.read_bytes()).hexdigest()
        document = DocumentRecord(
            document_id=document_id,
            course_id=course_id,
            document_title=doc_title,
            source_pdf=path_to_str(copied.resolve()),
            page_count=len(pages),
        )
        self.store.create_document(
            document,
            uploader_teacher_id=uploader_teacher_id,
            source_filename=source_filename or source_path.name,
            checksum_sha256=checksum,
            mime_type=DOCX_MIME,
        )
        self.store.create_pages(pages)
        chunks = self.chunker.chunk_pages(pages)
        self.store.upsert_chunks_for_document(course_id=course_id, document_id=document_id, chunks=chunks)
        self.store.update_document_status(document_id, status="ingested")
        logger.info("Ingested DOCX %s units for %s", len(pages), doc_title)
        return document

    def _ingest_pptx(
        self,
        source_path: Path,
        course_id: str,
        title: str | None,
        uploader_teacher_id: str | None,
        source_filename: str | None,
    ) -> DocumentRecord:
        try:
            from pptx import Presentation
        except Exception as exc:  # pragma: no cover - dependency may be missing
            raise RuntimeError("python-pptx is not installed. Add python-pptx to requirements.") from exc

        document_id, copied = self._copy_source(source_path, course_id, ".pptx")
        doc_title = title or source_path.stem
        prs = Presentation(copied)
        units: list[str] = []
        for slide in prs.slides:
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    value = str(shape.text).strip()
                    if value:
                        parts.append(value)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            slide_text = "\n".join(parts).strip()
            if notes:
                slide_text = f"{slide_text}\nNotes: {notes}".strip()
            units.append(slide_text if slide_text else f"Slide {len(units) + 1}")

        pages = self._build_text_units_as_pages(
            course_id=course_id,
            document_id=document_id,
            document_title=doc_title,
            units=units,
            unit_prefix="slide",
        )
        checksum = hashlib.sha256(copied.read_bytes()).hexdigest()
        document = DocumentRecord(
            document_id=document_id,
            course_id=course_id,
            document_title=doc_title,
            source_pdf=path_to_str(copied.resolve()),
            page_count=len(pages),
        )
        self.store.create_document(
            document,
            uploader_teacher_id=uploader_teacher_id,
            source_filename=source_filename or source_path.name,
            checksum_sha256=checksum,
            mime_type=PPTX_MIME,
        )
        self.store.create_pages(pages)
        chunks = self.chunker.chunk_pages(pages)
        self.store.upsert_chunks_for_document(course_id=course_id, document_id=document_id, chunks=chunks)
        self.store.update_document_status(document_id, status="ingested")
        logger.info("Ingested PPTX %s slides for %s", len(pages), doc_title)
        return document

    def _copy_source(self, source_path: Path, course_id: str, ext: str) -> tuple[str, Path]:
        document_id = str(uuid.uuid4())
        target = settings.documents_dir / course_id / f"{document_id}{ext}"
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source_path, target)
        return document_id, target

    @staticmethod
    def _build_text_units_as_pages(
        course_id: str,
        document_id: str,
        document_title: str,
        units: list[str],
        unit_prefix: str,
    ) -> list[PageRecord]:
        pages: list[PageRecord] = []
        for idx, text in enumerate(units, start=1):
            value = (text or "").strip()
            if not value:
                continue
            page = PageRecord(
                course_id=course_id,
                document_id=document_id,
                document_title=document_title,
                page_id=f"{document_id}_p{idx}",
                page_number=idx,
                image_path="",
                pdf_text_raw=value,
                ocr_text_raw="",
                ocr_text_clean=value,
                merged_text=value,
                text_source="pdf",
                pdf_text_quality=1.0,
                ocr_text_quality=1.0,
                language=detect_language(value),
                has_diagram=(unit_prefix == "slide"),
                has_table=False,
                has_code_like_text=False,
                has_large_image=(unit_prefix == "slide"),
            )
            pages.append(page)
        return pages

    @staticmethod
    def _build_video_pages(
        course_id: str,
        document_id: str,
        document_title: str,
        segments: list[TranscriptSegment],
    ) -> list[PageRecord]:
        pages: list[PageRecord] = []
        max_segment_seconds = max(20, settings.video_max_segment_seconds)
        bucket: list[TranscriptSegment] = []
        bucket_start: float | None = None
        bucket_end: float | None = None

        def flush_bucket() -> None:
            nonlocal bucket, bucket_start, bucket_end
            if not bucket or bucket_start is None or bucket_end is None:
                return
            text = " ".join(item.text for item in bucket).strip()
            idx = len(pages) + 1
            image_path = (
                f"video://{document_id}"
                f"?start={bucket_start:.3f}&end={bucket_end:.3f}&label={_format_time_range(bucket_start, bucket_end)}"
            )
            pages.append(
                PageRecord(
                    course_id=course_id,
                    document_id=document_id,
                    document_title=document_title,
                    page_id=f"{document_id}_p{idx}",
                    page_number=idx,
                    image_path=image_path,
                    pdf_text_raw=text,
                    ocr_text_raw="",
                    ocr_text_clean=text,
                    merged_text=text,
                    text_source="pdf",
                    pdf_text_quality=0.88,
                    ocr_text_quality=0.88,
                    language=detect_language(text),
                    has_diagram=False,
                    has_table=False,
                    has_code_like_text=False,
                    has_large_image=False,
                )
            )
            bucket = []
            bucket_start = None
            bucket_end = None

        for seg in segments:
            if bucket_start is None:
                bucket_start = seg.start_sec
                bucket_end = seg.end_sec
                bucket = [seg]
                continue
            next_end = seg.end_sec
            if (next_end - bucket_start) <= max_segment_seconds:
                bucket.append(seg)
                bucket_end = next_end
            else:
                flush_bucket()
                bucket_start = seg.start_sec
                bucket_end = seg.end_sec
                bucket = [seg]
        flush_bucket()
        return pages


def _mime_for_ext(ext: str) -> str:
    mapping = {
        ".pdf": "application/pdf",
        ".docx": DOCX_MIME,
        ".pptx": PPTX_MIME,
        ".mp4": "video/mp4",
        ".mov": "video/quicktime",
        ".m4v": "video/x-m4v",
        ".mkv": "video/x-matroska",
        ".webm": "video/webm",
    }
    return mapping.get(ext, "application/octet-stream")


def _format_time_range(start_sec: float, end_sec: float) -> str:
    return f"{_format_time(start_sec)}-{_format_time(end_sec)}"


def _format_time(value: float) -> str:
    total = int(max(0, value))
    h = total // 3600
    m = (total % 3600) // 60
    s = total % 60
    if h > 0:
        return f"{h:02d}:{m:02d}:{s:02d}"
    return f"{m:02d}:{s:02d}"
